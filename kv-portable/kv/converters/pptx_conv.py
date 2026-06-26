from __future__ import annotations

from pathlib import Path

from kv.converters.base import BaseConverter, ConvertResult


class PptxConverter(BaseConverter):
    source_type = "pptx"
    extensions = {".pptx"}

    def convert(self, path: Path) -> ConvertResult:
        text = self._extract(path)
        body = f"""# {path.stem}

## 원본
- 파일: `{path.name}`
- 유형: PPTX → 슬라이드 텍스트

{text}
"""
        return ConvertResult(title=path.stem, body=body, source_type=self.source_type)

    def _extract(self, path: Path) -> str:
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                lines: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(run.text for run in para.runs).strip()
                            if t:
                                lines.append(t)
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [c.text.strip() for c in row.cells]
                            lines.append("| " + " | ".join(cells) + " |")
                if lines:
                    parts.append(f"## 슬라이드 {i}\n\n" + "\n".join(lines))
            text = "\n\n".join(parts).strip()
            return text or "(슬라이드에서 텍스트를 찾지 못했습니다)"
        except ImportError:
            return "(PPTX 불가 — `pip install python-pptx` 필요)"
        except Exception as e:
            return f"(PPTX 변환 오류: {e})"
